#!/usr/bin/env python3
"""NS-climbing hackathon pitch deck -> .pptx (imports cleanly into Google Slides).
Brand matches the live site: near-black ground, vermilion->amber accents, wall imagery.
A dedicated QR slide points the room at the live cheer board (gasless, no wallet)."""
import os, segno
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE   = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "..", "assets")
OUT    = os.path.join(HERE, "NS_Climbing_Pitch.pptx")
QR_PNG = os.path.join(HERE, "cheer_qr.png")
CHEER_URL = "https://tinyurl.com/sendcheer"   # 301 -> cheer board (short = scans from the back row)

# ---- brand tokens (from the live site) ----
BG     = RGBColor(0x0E, 0x11, 0x16)   # near-black ground
PANEL  = RGBColor(0x16, 0x1B, 0x22)
INK    = RGBColor(0xF2, 0xF4, 0xF7)   # titles / statements
BODY   = RGBColor(0xC7, 0xCF, 0xDA)   # body
MUTE   = RGBColor(0x9A, 0xA4, 0xB2)   # captions
ACCENT = RGBColor(0xFF, 0x6B, 0x4A)   # vermilion
AMBER  = RGBColor(0xFF, 0xB2, 0x4A)   # amber
GREEN  = RGBColor(0x7F, 0xE0, 0xB2)   # "live" pill
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Arial"                       # universally present after Google Slides import

# ---- QR: dark modules on white, generous quiet zone, crisp for a projector ----
segno.make(CHEER_URL, error="h").save(QR_PNG, scale=24, border=3,
                                       dark="#0e1116", light="#ffffff")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s

def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return b, tf

def para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, first=False,
         spacing=1.0, tracking=None, before=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.line_spacing = spacing
    if before: p.space_before = Pt(before)
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
    return p, r

def bar(s, l, t, w=1.7, h=0.11, color=ACCENT):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
    b.shadow.inherit = False
    return b

def pill(s, l, t, text, color=GREEN, w=2.2):
    p = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.42))
    p.fill.solid(); p.fill.fore_color.rgb = BG
    p.line.color.rgb = color; p.line.width = Pt(1.25); p.shadow.inherit = False
    tf = p.text_frame; tf.word_wrap = False
    r = tf.paragraphs[0]; r.alignment = PP_ALIGN.CENTER
    run = r.add_run(); run.text = text
    run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = color; run.font.name = FONT
    return p

def _optimized(img, maxw=1800):
    """downscale a full-bleed image to a lean JPEG so the .pptx stays small."""
    from PIL import Image
    src = os.path.join(ASSETS, img)
    if not os.path.exists(src):
        return None
    out = os.path.join(HERE, "_opt_" + os.path.splitext(img)[0] + ".jpg")
    im = Image.open(src).convert("RGB")
    if im.width > maxw:
        im = im.resize((maxw, round(im.height * maxw / im.width)), Image.LANCZOS)
    im.save(out, "JPEG", quality=82, optimize=True)
    return out

def fullbleed(s, img, dark=0.62):
    """cover image + dark scrim so text stays legible."""
    path = _optimized(img)
    if path and os.path.exists(path):
        s.shapes.add_picture(path, 0, 0, width=SW, height=SH)
    ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    ov.fill.solid(); ov.fill.fore_color.rgb = BG
    ov.line.fill.background(); ov.shadow.inherit = False
    ov.fill.transparency = 0  # set below via alpha hack
    # python-pptx has no direct transparency; emulate with a solid near-black at given alpha
    _set_alpha(ov, int(dark * 255))
    return ov

def _set_alpha(shape, alpha):
    """add an <a:alpha> to a solidFill (0..255 -> 0..100000)."""
    from pptx.oxml.ns import qn
    sp = shape.fill._xPr.find(qn('a:solidFill'))
    srgb = sp.find(qn('a:srgbClr'))
    a = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha/255*100000))})
    srgb.append(a)

def footer(s):
    b, tf = tb(s, 0.55, 7.02, 8, 0.4)
    para(tf, "Network School  ·  send-climbing  ·  loringtonian.github.io/NS-climbing",
         11, MUTE, first=True, tracking=True)

def diagram(s, path, top=1.6, maxw=11.7, maxh=5.35):
    """center a generated architecture diagram, re-encoded lean (opaque dark
    image -> high-q JPEG at 1800px; labels stay crisp, deck stays small)."""
    from PIL import Image as _IM
    if not os.path.exists(path):
        return
    im = _IM.open(path).convert("RGB")
    if im.width > 1800:
        im = im.resize((1800, round(im.height * 1800 / im.width)), _IM.LANCZOS)
    path = os.path.splitext(path)[0] + "_opt.jpg"
    im.save(path, "JPEG", quality=90, optimize=True)
    iw, ih = _IM.open(path).size
    ar = iw / ih
    w = maxw; h = w / ar
    if h > maxh:
        h = maxh; w = h * ar
    s.shapes.add_picture(path, Inches((13.333 - w) / 2), Inches(top),
                         width=Inches(w), height=Inches(h))

# ============================ SLIDES ============================

def leadline(tf, tag, col, txt, size=20, first=False, before=13):
    """a colored bold lead word + body run on one paragraph (scannable)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(0 if first else before); p.line_spacing = 1.06
    a = p.add_run(); a.text = tag + "   "
    a.font.bold = True; a.font.size = Pt(size); a.font.color.rgb = col; a.font.name = FONT
    b = p.add_run(); b.text = txt
    b.font.size = Pt(size - 1.5); b.font.color.rgb = BODY; b.font.name = FONT

# 1 — TITLE
s = slide()
fullbleed(s, "bg_wall.jpg", dark=0.58)
_, tf = tb(s, 0.9, 2.35, 11.5, 3.2)
para(tf, "MEMBER INITIATIVE", 15, AMBER, bold=True, first=True)
para(tf, "We want a climbing wall", 54, INK, bold=True, before=10)
para(tf, "at Network School.", 54, ACCENT, bold=True)
para(tf, "A petition where the signatures are money.", 24, BODY, before=18)
footer(s)

# 2 — THE PROBLEM = MEASUREMENT
s = slide()
bar(s, 0.9, 1.1)
_, tf = tb(s, 0.9, 2.0, 11.6, 3.8)
para(tf, "Collective action dies at measurement.", 44, INK, bold=True, first=True)
para(tf, "Last year a petition almost flipped it. Almost.", 24, BODY, before=22)
para(tf, "Nobody knows how many people truly want the thing — so no one goes first. Ask for "
        "money too early and you kill the signal. Ask for nothing, and the signal is worth "
        "nothing.", 22, MUTE, before=10, spacing=1.1)
footer(s)

# 3 — THE THESIS: TWO LAYERS
s = slide()
_, tf = tb(s, 0.9, 0.95, 11.6, 1.0)
para(tf, "So we split it in two — priced on purpose.", 38, INK, bold=True, first=True)
_, t2 = tb(s, 0.9, 2.5, 11.6, 2.1)
leadline(t2, "FREE", AMBER, "smash a button. Desire, measured — and on-chain.", size=24, first=True)
leadline(t2, "LOCKED", ACCENT, "USDC in a pool no one can empty. Commitment, permanent.", size=24, before=16)
_, t3 = tb(s, 0.9, 4.95, 11.6, 1.9)
para(t3, "The rollup is where a community discovers it wants something.", 26, INK, bold=True, first=True, spacing=1.12)
para(t3, "The mainnet contract is where it proves it.", 26, AMBER, bold=True, spacing=1.12)
footer(s)

# 4 — LAYER 1: THE CHEER BOARD (MagicBlock, emphasized)
s = slide()
pill(s, 0.9, 1.0, "●  LAYER 1 · MAGICBLOCK EPHEMERAL ROLLUP", AMBER, w=5.7)
_, tf = tb(s, 0.9, 1.9, 11.6, 3.2)
para(tf, "Saying “yes” should cost nothing.", 40, INK, bold=True, first=True)
para(tf, "The cheer board is free, gasless, wallet-less — you just smash. That is only possible "
        "on an Ephemeral Rollup; on base-layer Solana, a wallet and a fee per tap would have "
        "killed it dead.", 21, BODY, before=18, spacing=1.1)
_, t2 = tb(s, 0.9, 5.15, 11.6, 1.3)
para(t2, "25,000+ cheers · 30 people", 38, AMBER, bold=True, first=True)
para(t2, "live, in the room, on-chain — not a Google Form anyone can stuff.", 19, MUTE, before=4)
footer(s)

# 5 — HOW THE ROLLUP WORKS (lifecycle)
s = slide()
pill(s, 0.9, 1.0, "●  HOW THE ROLLUP WORKS", AMBER, w=3.9)
_, tf = tb(s, 0.9, 1.75, 11.6, 1.5)
para(tf, "On-demand SVM runtimes on Solana.", 36, INK, bold=True, first=True)
para(tf, "Delegate an account to the rollup, run it at ~10 ms and gasless, then commit state "
        "back to the base layer.", 18, MUTE, before=8, spacing=1.06)
_, t2 = tb(s, 0.9, 3.5, 11.6, 2.7)
leadline(t2, "DELEGATE", ACCENT, "we hand the rollup one counter PDA — a single number, never funds.", first=True)
leadline(t2, "CHEER", AMBER, "every tap is its own gasless transaction, ~10 ms, signed by a throwaway "
         "in-page key. No wallet, no SOL. ~366 / sec.")
leadline(t2, "COMMIT", GREEN, "undelegate, and the final tally lands on Solana L1 as a permanent record.")
_, t3 = tb(s, 0.9, 6.35, 11.6, 0.7)
para(t3, "Tallied straight from the chain, not a database — no server. Your face is your "
        "pubkey; no name ever touches the chain.", 15, MUTE, first=True)
footer(s)

# 6 — MAGICBLOCK ARCHITECTURE DIAGRAM
s = slide()
_, tf = tb(s, 0.9, 0.5, 11.6, 0.95)
para(tf, "One PDA — delegated to a rollup, committed back.", 30, INK, bold=True, first=True)
para(tf, "The shape of the cheer program.", 17, MUTE, before=6)
diagram(s, os.path.join(HERE, "arch_magicblock.png"), top=1.78)
footer(s)

# 7 — LAYER 2: THE ESCROW (turn + three exits)
s = slide()
pill(s, 0.9, 1.0, "●  LAYER 2 · LIVE ON SOLANA MAINNET", GREEN, w=4.8)
_, tf = tb(s, 0.9, 1.85, 11.6, 4.6)
para(tf, "Then — and only then — the money.", 40, INK, bold=True, first=True)
para(tf, "USDC locks into a pool no person can empty. Out only together, weighted by dollars:",
     22, BODY, before=16)
para(tf, "BUILD — a dollar-majority backs the organizer's address → funds the wall.", 20, AMBER, bold=True, before=13)
para(tf, "DISSOLVE — a dollar-majority votes no-confidence → everyone refunded.", 20, ACCENT, bold=True, before=7)
para(tf, "TIMEOUT — 180 days pass → everyone refunded, automatically.", 20, GREEN, bold=True, before=7)
para(tf, "No one — not the organizer, not anyone — can move a cent alone.", 21, INK, bold=True, before=15)
footer(s)

# 8 — ESCROW: REAL + IMMUTABLE (specs)
s = slide()
_, tf = tb(s, 0.9, 0.9, 11.6, 4.9)
para(tf, "Real, and immutable.", 40, INK, bold=True, first=True)
para(tf, "An Anchor program on Solana — the upgrade authority is burned, so the bytecode can "
        "never change, not even by us.", 21, BODY, before=18, spacing=1.08)
para(tf, "USDC locks in a program vault PDA. Each depositor gets a non-transferable receipt "
        "PDA — badge, vote, and refund in one.", 21, BODY, before=8, spacing=1.08)
para(tf, "Release requires  payout_vote_amount × 2 > total_escrowed  — a majority of dollars, "
        "not wallets. Refunds are permissionless; funds only ever return to each depositor.",
     21, BODY, before=8, spacing=1.08)
para(tf, "Program  2PAg6iMEzPQnfzVmKdeUDctmmCYwts46Y5GEZBUDA4KJ", 14, MUTE, before=16)
footer(s)

# 9 — ESCROW ARCHITECTURE DIAGRAM
s = slide()
_, tf = tb(s, 0.9, 0.5, 11.6, 0.95)
para(tf, "The shape of the contract.", 30, INK, bold=True, first=True)
para(tf, "Deposits in. Three collective ways out. No individual withdraw.", 17, MUTE, before=6)
diagram(s, os.path.join(HERE, "arch_escrow.png"), top=1.78)
footer(s)

# 10 — SCAN TO CHEER (the live moment)
s = slide()
_, tf = tb(s, 0.85, 1.15, 6.6, 5.0, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Scan to cheer.", 52, INK, bold=True, first=True)
para(tf, "Right now, from your seat.", 30, AMBER, bold=True, before=6)
para(tf, "No wallet. No cost. Tap as many times as you mean it — "
        "every tap is its own gasless transaction on the rollup.", 22, BODY, before=24, spacing=1.12)
para(tf, "tinyurl.com/sendcheer", 22, MUTE, bold=True, before=20)
qx, qy, qs = 8.15, 1.55, 4.35
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(qx-0.28), Inches(qy-0.28),
                           Inches(qs+0.56), Inches(qs+0.56))
panel.fill.solid(); panel.fill.fore_color.rgb = WHITE; panel.line.fill.background()
panel.shadow.inherit = False
s.shapes.add_picture(QR_PNG, Inches(qx), Inches(qy), width=Inches(qs), height=Inches(qs))
footer(s)

# 11 — WHY IT'S FUNDABLE (honest, human-driven)
s = slide()
bar(s, 0.9, 1.1)
_, tf = tb(s, 0.9, 1.5, 11.6, 5.0)
para(tf, "Priced demand, not a signature list.", 38, INK, bold=True, first=True)
para(tf, "This is Curious's own fundable idea #25 — communities set a goal, pool resources, and "
        "execute together — shipped, live, and immutable.", 21, BODY, before=20, spacing=1.08)
para(tf, "Human-driven, honestly: people cheer, people deposit, people vote. No autonomous agent "
        "DAO here — the agent layer is the honest next step, and agents.md already lets any AI "
        "audit the contract (and could place a deposit).", 19, MUTE, before=10, spacing=1.08)
para(tf, "A live wedge, not a slide: a real contract, real cheers, real on-chain usage this weekend.",
     20, AMBER, bold=True, before=12)
footer(s)

# 12 — CLOSE (ask your agent + build it)
s = slide()
fullbleed(s, "gym_concept_v0.png", dark=0.66)
_, tf = tb(s, 0.9, 1.55, 8.6, 4.4)
para(tf, "Don't trust us. Ask your agent.", 38, INK, bold=True, first=True)
para(tf, "The code is public and the contract is immutable. VERIFY_IT.md has a prompt you paste "
        "into your own AI to audit every claim against the chain — right now.", 18, BODY, before=14, spacing=1.1)
para(tf, "A signature says “sure, whatever.” Locked money says “build it.”", 26, AMBER, bold=True, before=22, spacing=1.14)
para(tf, "tinyurl.com/sendcheer", 20, INK, bold=True, before=14)
qs2 = 2.0
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.9), Inches(4.8),
                           Inches(qs2+0.3), Inches(qs2+0.3))
panel.fill.solid(); panel.fill.fore_color.rgb = WHITE; panel.line.fill.background()
panel.shadow.inherit = False
s.shapes.add_picture(QR_PNG, Inches(11.05), Inches(4.95), width=Inches(qs2), height=Inches(qs2))
footer(s)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides._sldIdLst))
