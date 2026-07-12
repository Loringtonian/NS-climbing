#!/usr/bin/env python3
"""One extra slide: the cross-chain / embedded-wallet (Privy) deposit path.

Diagram-led, almost no text — the brand marks do the talking. Standalone so the
result can be pasted straight into the Google Slides deck. Same brand tokens as
build_pitch.py. Emits a 1-slide .pptx and a PNG render.

Says nothing about any staging URL: this describes the mechanism, not a page.
Regenerate the graphic with gen_anychain.py (nano-banana-pro).
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE    = os.path.dirname(os.path.abspath(__file__))
OUT     = os.path.join(HERE, "anychain_slide.pptx")
DIAGRAM = os.path.join(HERE, "arch_anychain.png")
CROPPED = os.path.join(HERE, "arch_anychain_crop.jpg")

BG     = RGBColor(0x0E, 0x11, 0x16)
INK    = RGBColor(0xF2, 0xF4, 0xF7)
MUTE   = RGBColor(0x9A, 0xA4, 0xB2)
AMBER  = RGBColor(0xFF, 0xB2, 0x4A)
FONT   = "Arial"

SW, SH = Inches(13.333), Inches(7.5)


def crop_to_content(src, dst, pad=55, thresh=80, min_px=6):
    """The generator leaves huge dead space above/below the flow. Trim to the band
    that actually has content. A row counts as content only if it has several genuinely
    bright pixels — a plain max() gets fooled by the soft glow bleeding into the ground,
    which trims almost nothing and leaves the graphic too tall for the slide."""
    im = Image.open(src).convert("RGB")
    g = im.convert("L")
    w, h = g.size
    rows = [y for y in range(h)
            if sum(1 for v in g.crop((0, y, w, y + 1)).getdata() if v > thresh) >= min_px]
    top, bot = (max(0, rows[0] - pad), min(h, rows[-1] + pad)) if rows else (0, h)
    im = im.crop((0, top, w, bot))
    im.save(dst, "JPEG", quality=92, optimize=True)
    return im.size


dw, dh = crop_to_content(DIAGRAM, CROPPED)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG


def tb(l, t, w, h):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    return tf


def para(tf, text, size, color, bold=False, first=False, spacing=1.0, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.line_spacing = spacing
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT


pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.62), Inches(6.2), Inches(0.42))
pill.fill.solid(); pill.fill.fore_color.rgb = BG
pill.line.color.rgb = AMBER; pill.line.width = Pt(1.25); pill.shadow.inherit = False
r = pill.text_frame.paragraphs[0]; r.alignment = PP_ALIGN.CENTER
run = r.add_run(); run.text = "●  ANY CHAIN · ANY WALLET · NO WALLET AT ALL"
run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = AMBER; run.font.name = FONT

tf = tb(0.9, 1.22, 11.9, 0.7)
para(tf, "Bring money from anywhere. The contract never moves.", 28, INK, bold=True, first=True,
     spacing=1.04)

# the graphic runs full-bleed: the generated ground is the same near-black as the slide,
# so edge-to-edge reads as one surface instead of a pasted-in box.
gw = 13.333
gh = gw * dh / dw
gtop = 2.10
s.shapes.add_picture(CROPPED, 0, Inches(gtop), width=Inches(gw), height=Inches(gh))

# one honest line. The graphic already said the rest.
tf = tb(0.9, gtop + gh + 0.38, 11.9, 0.9)
para(tf, "The escrow stays trustless. Privy and the relayer are plumbing around it — "
         "they can help money in; they can never move a cent out.",
     17, AMBER, bold=True, first=True, spacing=1.1)

tf = tb(0.55, 7.02, 8, 0.4)
para(tf, "Network School  ·  send-climbing  ·  loringtonian.github.io/NS-climbing", 11, MUTE, first=True)

prs.save(OUT)
print(f"cropped graphic: {dw}x{dh} -> {gw:.1f}x{gh:.2f} in")
print("saved:", OUT)
